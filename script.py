#-*- coding:utf-8 -*-

import sys
import json
import time
import glob, os
import itertools as IT
import zipfile
from shutil import copyfile
from shutil import rmtree
from datetime import datetime
import numpy as np

# External py file
import config
import connection
from connection import *



SEVERITY = ['Critical',
            'High',
            'Medium',
            'Low',
			'Info']

try:
    from docxtpl import DocxTemplate, Listing, InlineImage, RichText
except ImportError:
    print("Missing docxtpl library. Please install it using PIP. https://pypi.org/project/docxtpl/")
    exit()

try:
    from docx.shared import Mm
except ImportError:
    print("Missing python-docx library. Please install it using PIP. https://pypi.org/project/python-docx/")
    exit()

try:
    import openpyxl
except ImportError:
    print("Missing openpyxl library. Please install it using PIP. https://pypi.org/project/openpyxl/")
    exit()

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Missing python-pptx library. Please install it using PIP. https://pypi.org/project/python-pptx/")
    exit()

from PIL import Image
 
os.chdir("/home/dojo/templates")

#This fuction groups normal Python arrays in groups of N, where N is a number.
#If N = 2 the output should be:
#[['5', '7'], ['8', '9'], ['4', '3'], ['1', '7'], ['2', '3'], ['1', '3'], ['0', '3'], ['7', '3'], ['4', '3'], ['6', '2'], ['2']]
#If Endpoint_Count < 10 then N = X, Else If Endpoint_Count < 30 then N = Y...
def grouper(n, iterable):
    iterable = iter(iterable)
    return iter(lambda: list(IT.islice(iterable, n)), [])


def generate_doc(all_endpoints, finding_endpoints, critical_findings, high_findings, medium_findings, low_findings, count_critical_findings, count_high_findings, count_medium_findings, count_low_findings, total_findings, list_finding_images, client_name, project_name, leader_name, p_start_date, p_end_date, path_project):
    doc = DocxTemplate("/home/dojo/templates/template.docx")
    os.chdir("/home/dojo/templates")
    for i in range(len(list_finding_images)):
        path_image = list_finding_images[i]['image']
        path_image = InlineImage(doc, '/home/dojo/media/{}'.format(path_image), width=Mm(150), height=Mm(60))
        list_finding_images[i]['image'] = path_image


    for i in range(len(critical_findings)):
        critical_findings[i]['description'] = critical_findings[i]['description'].replace('\n', '\a ')
        critical_findings[i]['description'] = critical_findings[i]['description'].replace('```', '')
        critical_findings[i]['description'] = RichText('{}'.format(critical_findings[i]['description']))

        critical_findings[i]['impact'] = critical_findings[i]['impact'].replace('\n', '\a ')
        critical_findings[i]['impact'] = RichText('{}'.format(critical_findings[i]['impact']))

        critical_findings[i]['refs'] = critical_findings[i]['refs'].replace('\n', '\a ')
        critical_findings[i]['refs'] = RichText('{}'.format(critical_findings[i]['refs']))
        ###################################################################################################
        high_findings[i]['description'] = high_findings[i]['description'].replace('\n', '\a ')
        high_findings[i]['description'] = high_findings[i]['description'].replace('```', '')
        high_findings[i]['description'] = RichText('{}'.format(high_findings[i]['description']))

        high_findings[i]['impact'] = high_findings[i]['impact'].replace('\n', '\a ')
        high_findings[i]['impact'] = RichText('{}'.format(high_findings[i]['impact']))

        high_findings[i]['refs'] = high_findings[i]['refs'].replace('\n', '\a ')
        high_findings[i]['refs'] = RichText('{}'.format(high_findings[i]['refs']))
        ###################################################################################################
        medium_findings[i]['description'] = medium_findings[i]['description'].replace('\n', '\a ')
        medium_findings[i]['description'] = medium_findings[i]['description'].replace('```', '')
        medium_findings[i]['description'] = RichText('{}'.format(medium_findings[i]['description']))

        medium_findings[i]['impact'] = medium_findings[i]['impact'].replace('\n', '\a ')
        medium_findings[i]['impact'] = RichText('{}'.format(medium_findings[i]['impact']))

        medium_findings[i]['refs'] = medium_findings[i]['refs'].replace('\n', '\a ')
        medium_findings[i]['refs'] = RichText('{}'.format(medium_findings[i]['refs']))
        ###################################################################################################
        low_findings[i]['description'] = low_findings[i]['description'].replace('\n', '\a ')
        low_findings[i]['description'] = low_findings[i]['description'].replace('```', '')
        low_findings[i]['description'] = RichText('{}'.format(low_findings[i]['description']))

        low_findings[i]['impact'] = low_findings[i]['impact'].replace('\n', '\a ')
        low_findings[i]['impact'] = RichText('{}'.format(low_findings[i]['impact']))

        low_findings[i]['refs'] = low_findings[i]['refs'].replace('\n', '\a ')
        low_findings[i]['refs'] = RichText('{}'.format(low_findings[i]['refs']))
        ###################################################################################################

    context = {
        'finding_images' : list_finding_images,

        'client_name': client_name,
        'project_name': project_name,
        'leader_name': leader_name, 
        
        'p_start_date': p_start_date,
        'p_end_date': p_end_date,
        
        'all_endpoints': all_endpoints,
        'finding_endpoints': finding_endpoints, 
        
        'low_findings': low_findings,
        'count_low_findings': count_low_findings,
        
        'medium_findings': medium_findings,
        'count_medium_findings': count_medium_findings,
        
        'high_findings': high_findings,
        'count_high_findings': count_high_findings,
        
        'critical_findings': critical_findings,
        'count_critical_findings': count_critical_findings,
        
        'total_findings': total_findings,
    }

    doc.render(context, autoescape=True)
    doc.save(path_project + "/" + client_name + "_" + project_name + ".docx")


def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]
 
    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size
 
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width
 
    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)
 
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
 
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side


def generate_pptx(count_critical_findings, count_high_findings, count_medium_findings, count_low_findings, path_project, client_name, project_name, critical_findings, list_finding_images):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Critical Issues"

    os.chdir("/home/dojo/media/")
    for i in range(len(critical_findings)):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = critical_findings[i]['title']

        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tfa = txBox.text_frame

        tfa = body_shape.text_frame
        tfa.text = critical_findings[i]['mitigation']

        p = tfa.add_paragraph()
        p.text = critical_findings[i]['impact']

        for j in range(len(list_finding_images)):
            if(critical_findings[i]['finding_id'] == list_finding_images[j]['finding_id']):
                img_path = list_finding_images[j]['image']

                picture_with_caption = prs.slide_layouts[8]
                slide = prs.slides.add_slide(picture_with_caption)

                _add_image(slide,1,img_path)
                #image = slide.placeholders[1].insert_picture(img_path)
                title = slide.shapes.title.text = critical_findings[i]['title']
                caption = slide.placeholders[2].text = list_finding_images[j]['caption']
                #left = top = (prs.slide_width - image.width) / 2
                #pic = slide.shapes.add_picture(img_path, left, top)

    prs.save(path_project + '/' + client_name + '_' + project_name + '.pptx')


def validate_ids(client_id, project_id):

    sql_validate_ids = ("""
        SELECT EXISTS(SELECT * FROM dojo_engagement WHERE product_id = '%s' and id = '%s') AS client_project_exists
        """ % (client_id, project_id))
    
    client_project_exists = select_sql(sql_validate_ids)

    return int(client_project_exists['data'][0][0])


def check_scans_existence(project_id):
    #Here we can assume that we always have a scan name and simple concat "Scan Name (Scan Type)" - "Initial Scan (Nessus Scan)"
    #WHY ARE WE CONCATING A TITLE? NO NEED
    sql_check_scans_existence_query = ("""
        SELECT 
            EXISTS(SELECT dojo_test.id,
            concat( if(title is not null or title = '', concat(dojo_test.title,'(',dojo_test_type.name,')' ), dojo_test_type.name) ) AS scan_name 
        FROM dojo_test 
            INNER JOIN dojo_engagement 
                ON dojo_test.engagement_id = dojo_engagement.id 
            INNER JOIN dojo_test_type 
                ON dojo_test_type.id = test_type_id 
        WHERE dojo_engagement.id = '%s') AS project_scan_exists
        """ % (project_id))
    scans_existence = select_sql(sql_check_scans_existence_query)

    return int(scans_existence['data'][0][0])


def get_project_info(project_id):
    project_info_query = ("""
        SELECT 
            dojo_product.name AS client_name, 
            dojo_engagement.name AS project_name, 
            CONCAT(auth_user.first_name, ' ', auth_user.last_name) as leader_name, 
            target_start, 
            target_end 
    FROM auth_user 
        INNER JOIN dojo_engagement 
            ON auth_user.id = dojo_engagement.lead_id 
        INNER JOIN dojo_product 
            ON dojo_product.id = dojo_engagement.product_id 
        WHERE dojo_engagement.id = '%s'
    """ % (project_id))

    tuple_project_info = select_sql(project_info_query)
    
    row_headers = tuple_project_info['row_headers']
    records_project_info = tuple_project_info['data']
    
    data_project_info = []
    for data in records_project_info:
        data_project_info.append(dict(zip(row_headers,data)))

        project_info = json.dumps(data_project_info, indent=4, sort_keys=True, default=str)
        project_info = json.loads(project_info)

    return project_info


def get_all_endpoints(project_id):
    sql_select_all_endpoints_query = ("""
        SELECT 
            DISTINCT SUBSTRING_INDEX(host, ':', 1)
        FROM dojo_endpoint
            INNER JOIN dojo_product
                ON dojo_product.id = dojo_endpoint.product_id
            INNER JOIN dojo_engagement 
                ON dojo_engagement.product_id  = dojo_product.id
        WHERE dojo_engagement.id = '%s'
            ORDER BY SUBSTRING_INDEX(host, ':', 1) ASC
    """ % (project_id))

    tuple_all_endpoints = select_sql(sql_select_all_endpoints_query)

    records_all_endpoints = tuple_all_endpoints['data']
    endpoint_length = tuple_all_endpoints['row_count']

    endpoint_array = [None] * endpoint_length
    endpoint_array[0] = 0
    number_endpoint = 0

    for row in records_all_endpoints:
       endpoint_array[number_endpoint] = row[0]
       number_endpoint += 1

    return endpoint_array


def get_finding_endpoints(project_id):
    sql_select_finding_endpoints_query = ("""
        SELECT 
            dojo_endpoint.id, 
            dojo_endpoint.host, 
            dojo_endpoint.protocol, 
            dojo_finding_endpoints.finding_id 
        FROM dojo_endpoint 
            INNER JOIN dojo_finding_endpoints 
                ON endpoint_id = dojo_endpoint.id 
            INNER JOIN dojo_finding 
                ON finding_id = dojo_finding.id 
            INNER join dojo_test 
                ON dojo_test.id = dojo_finding.test_id 
            INNER JOIN dojo_engagement 
                ON dojo_engagement.id = dojo_test.engagement_id 
        WHERE engagement_id = '%s'
            ORDER BY finding_id ASC
    """ % (project_id))

    tuple_finding_endpoints = select_sql(sql_select_finding_endpoints_query)
    
    row_headers_finding_endpoints = tuple_finding_endpoints['row_headers']
    records_finding_endpoints = tuple_finding_endpoints['data']

    data_finding_endpoints = []
    for data in records_finding_endpoints:
        data_finding_endpoints.append(dict(zip(row_headers_finding_endpoints,data)))

        finding_endpoints = json.dumps(data_finding_endpoints, indent=4, sort_keys=True, default=str)
        finding_endpoints = json.loads(finding_endpoints)

    return finding_endpoints


def get_critical_findings(project_id):
    sql_select_critical_findings_query = ("""
        SELECT dojo_finding.id AS finding_id, 
            dojo_finding.title, 
            date, 
            cwe, 
            url, 
            severity, 
            dojo_finding.description, 
            mitigation, 
            impact, 
            active, 
            dojo_finding.created, 
            cve, 
            refs,
            CONCAT(auth_user.first_name, ' ', auth_user.last_name) as reporter_name
        FROM dojo_finding 
            INNER JOIN dojo_test 
                ON dojo_test.id = dojo_finding.test_id
            INNER JOIN auth_user
                ON auth_user.id = dojo_finding.reporter_id
        WHERE dojo_test.engagement_id = '%s'
            AND severity='Critical' 
            AND dojo_finding.false_p = 0
    """ % (project_id))

    critical_findings = select_sql(sql_select_critical_findings_query)
    
    return critical_findings


def get_high_findings(project_id):
    sql_select_high_findings_query = ("""
        SELECT 
            dojo_finding.id AS finding_id, 
            dojo_finding.title, 
            date, 
            cwe, 
            url, 
            severity, 
            dojo_finding.description, 
            mitigation, 
            impact, 
            active, 
            dojo_finding.created, 
            cve, 
            refs,
            CONCAT(auth_user.first_name, ' ', auth_user.last_name) as reporter_name
        FROM dojo_finding 
            INNER JOIN dojo_test 
                ON dojo_test.id = dojo_finding.test_id
            INNER JOIN auth_user
                ON auth_user.id = dojo_finding.reporter_id
        WHERE dojo_test.engagement_id = '%s'
            AND severity='High' 
            AND dojo_finding.false_p = 0
    """ % (project_id))

    high_findings = select_sql(sql_select_high_findings_query)
    
    return high_findings


def get_medium_findings(project_id):
    sql_select_medium_findings_query = ("""
        SELECT 
            dojo_finding.id AS finding_id,
            dojo_finding.title,
            date,
            cwe, 
            url, 
            severity,
            dojo_finding.description, 
            mitigation, 
            impact, 
            active, 
            dojo_finding.created, 
            cve, 
            refs,
            CONCAT(auth_user.first_name, ' ', auth_user.last_name) as reporter_name
        FROM dojo_finding 
            INNER JOIN dojo_test 
                ON dojo_test.id = dojo_finding.test_id
            INNER JOIN auth_user
                ON auth_user.id = dojo_finding.reporter_id
        WHERE dojo_test.engagement_id = '%s'
            AND severity='Medium' 
            AND dojo_finding.false_p = 0
    """ % (project_id))

    medium_findings = select_sql(sql_select_medium_findings_query)
    
    return medium_findings


def get_low_findings(project_id):
    sql_select_low_findings_query = ("""
        SELECT 
            dojo_finding.id AS finding_id,
            dojo_finding.title,
            date, 
            cwe, 
            url, 
            severity, 
            dojo_finding.description, 
            mitigation, 
            impact, 
            active, 
            dojo_finding.created, 
            cve, 
            refs,
            CONCAT(auth_user.first_name, ' ', auth_user.last_name) as reporter_name
        FROM dojo_finding 
            INNER JOIN dojo_test 
                ON dojo_test.id = dojo_finding.test_id
            INNER JOIN auth_user
                ON auth_user.id = dojo_finding.reporter_id
        WHERE dojo_test.engagement_id = '%s' 
            AND severity='Low' 
            AND dojo_finding.false_p = 0
    """ % (project_id))

    low_findings = select_sql(sql_select_low_findings_query)
    
    return low_findings


def get_finding_images(project_id):
    sql_select_finding_image_query = ("""
    SELECT 
        dojo_finding.id AS finding_id, 
        dojo_findingimage.image,
        dojo_findingimage.caption 
        FROM dojo_finding 
            INNER JOIN dojo_finding_images 
                ON dojo_finding_images.finding_id = dojo_finding.id 
            INNER JOIN dojo_findingimage 
                ON dojo_findingimage.id = dojo_finding_images.findingimage_id 
            INNER JOIN dojo_test 
                ON dojo_test.id = dojo_finding.test_id 
            INNER JOIN dojo_engagement 
                ON dojo_engagement.id = dojo_test.engagement_id 
        WHERE dojo_engagement.id = '%s'
        AND dojo_finding.active = 1
    """ % (project_id))

    finding_images = select_sql(sql_select_finding_image_query)
    
    return finding_images


def get_all_info(project_id):
    project_info = get_project_info(project_id)
    
    client_name = project_info[0]['client_name']
    project_name = project_info[0]['project_name']
    leader_name = project_info[0]['leader_name']
    
    p_start_date = project_info[0]['target_start']
    p_end_date = project_info[0]['target_end']
    
    all_endpoints = get_all_endpoints(project_id)
    finding_endpoints = get_finding_endpoints(project_id)


    tuple_critical_findings = get_critical_findings(project_id)
    
    count_critical_findings = tuple_critical_findings['row_count']
    records_critical_findings = tuple_critical_findings['data']
    row_headers_critical_findings = tuple_critical_findings['row_headers']
    
    if count_critical_findings == 0:
        critical_findings = "No Critical Risk Issues found."
    else:
        list_critical_findings = []
        critical_findings = []
        for data_critical_findings in records_critical_findings:
            #list_critical_findings.append(dict(zip(row_headers_critical_findings, data_critical_findings)))
            critical_findings.append(dict(zip(row_headers_critical_findings, data_critical_findings)))

        #critical_findings = json.dumps(list_critical_findings, indent=4, sort_keys=True, default=str)
        #critical_findings = json.loads(critical_findings)


    tuple_high_findings = get_high_findings(project_id)
    
    count_high_findings = tuple_high_findings['row_count']
    records_high_findings = tuple_high_findings['data']
    row_headers_high_findings = tuple_high_findings['row_headers']
    
    if count_high_findings == 0:
        high_findings = "No High Risk Issues found."
    else:
        list_high_findings = []
        high_findings = []
        for data_high_findings in records_high_findings:
            #list_high_findings.append(dict(zip(row_headers_high_findings,data_high_findings)))
            high_findings.append(dict(zip(row_headers_high_findings,data_high_findings)))

        #high_findings = json.dumps(list_high_findings, indent=4, sort_keys=True, default=str)
        #high_findings = json.loads(high_findings)
    

    tuple_medium_findings = get_medium_findings(project_id)
    
    count_medium_findings = tuple_medium_findings['row_count']
    records_medium_findings = tuple_medium_findings['data']
    row_headers_medium_findings = tuple_medium_findings['row_headers']
    
    if count_medium_findings == 0:
        medium_findings = "No Medium Risk Issues found."
    else:
        list_medium_findings = []
        medium_findings = []
        for data_medium_findings in records_medium_findings:
            #list_medium_findings.append(dict(zip(row_headers_medium_findings,data_medium_findings)))
            medium_findings.append(dict(zip(row_headers_medium_findings,data_medium_findings)))

        #medium_findings = json.dumps(list_medium_findings, indent=4, sort_keys=True, default=str)
        #medium_findings = json.loads(medium_findings)

    
    tuple_low_findings = get_low_findings(project_id)
    
    count_low_findings = tuple_low_findings['row_count']
    records_low_findings = tuple_low_findings['data']
    row_headers_low_findings = tuple_low_findings['row_headers']
    
    if count_low_findings == 0:
        low_findings = "No Low Risk Issues found."
    else:
        list_low_findings = []
        low_findings = []
        for data_low_findings in records_low_findings:
            #list_low_findings.append(dict(zip(row_headers_low_findings,data_low_findings)))
            low_findings.append(dict(zip(row_headers_low_findings,data_low_findings)))

        #low_findings = json.dumps(list_low_findings, indent=4, sort_keys=True, default=str)
        #low_findings = json.loads(low_findings)
 
    tuple_finding_images = get_finding_images(project_id)
    count_finding_images = tuple_finding_images['row_count']
    records_finding_images = tuple_finding_images['data']
    row_headers_images = tuple_finding_images['row_headers']
    
    list_finding_images = []
    for images in records_finding_images:
        list_finding_images.append(dict(zip(row_headers_images,images)))

    path_client = "/home/dojo/clients/" + client_name + ""

    try:
        os.mkdir(path_client)
    except OSError:
        print ("Creation of the directory %s failed" % path_client)
    else:
        print ("Successfully created the directory %s " % path_client)

        
    path_project = path_client + "/" + project_name + ""

    try:
        os.mkdir(path_project)
    except OSError:
        print ("Creation of the directory %s failed" % path_project)
    else:
        print ("Successfully created the directory %s " % path_project)


    clientProject = client_name + "_" + project_name
    create_graph(count_critical_findings, count_high_findings, count_medium_findings, count_low_findings, path_project, clientProject)
    total_findings = count_critical_findings + count_high_findings + count_medium_findings + count_low_findings

    generate_pptx(count_critical_findings, count_high_findings, count_medium_findings, count_low_findings, path_project, client_name, project_name, critical_findings, list_finding_images)
    generate_doc(all_endpoints, finding_endpoints, critical_findings, high_findings, medium_findings, low_findings, count_critical_findings, count_high_findings, count_medium_findings, count_low_findings, total_findings, list_finding_images, client_name, project_name, leader_name, p_start_date, p_end_date, path_project)
    

def create_graph(count_critical_findings, count_high_findings, count_medium_findings, count_low_findings, path_project, clientProject):    
    workbook = openpyxl.load_workbook(filename='/home/dojo/templates/chart_template.xlsm', read_only=False, keep_vba=True)
    worksheet = workbook.get_sheet_by_name('Charts')

    worksheet['C6'] = count_low_findings
    worksheet['C7'] = count_medium_findings
    worksheet['C8'] = count_high_findings
    worksheet['C9'] = count_critical_findings

    worksheet['C15'] = count_low_findings
    worksheet['C16'] = count_medium_findings
    worksheet['C17'] = count_high_findings
    worksheet['C18'] = count_critical_findings

    workbook.save(path_project + '/' + clientProject + '.xlsx')

    PAD = os.getcwd()

    with zipfile.ZipFile('/home/dojo/templates/chart_template.xlsm', 'r') as z:
        z.extractall('./xlsm/')

    with zipfile.ZipFile(path_project + '/' + clientProject + '.xlsx', 'r') as z:
        z.extractall('./xlsx/')

    copyfile('./xlsm/xl/drawings/drawing1.xml','./xlsx/xl/drawings/drawing1.xml')

    z = zipfile.ZipFile(path_project + '/' + clientProject + '.zip', 'w')

    os.chdir('./xlsx')

    for root, dirs, files in os.walk('./'):
            for file in files:
                z.write(os.path.join(root, file))
    z.close()

    os.chdir(PAD)
    rmtree('./xlsm/')
    rmtree('./xlsx/')
    os.remove(path_project + '/' + clientProject + '.xlsx')
    os.rename(path_project + '/' + clientProject + '.zip', path_project + '/' + clientProject + '.xlsm')

  
def main():
    try:
        if len(sys.argv) > 2:
            client_id = sys.argv[1]
            project_id = sys.argv[2]
            #template = sys.argv[3]
            try:  
                if(validate_ids(client_id,project_id) == 1):
                    #print("validate_ids(true)")
                    if(check_scans_existence(project_id) == 1):
                        get_all_info(project_id)
                    else:
                        None
                else:
                    None

            except Exception as inst:
                print('     [!] Error processing file:\n' + str(inst))
        else:
            print("\n     Usage: ./run.py \"client ID\" \"project ID\" \"template name\"")

    except Exception as inst:
        print('     Critical Error: '+ str(inst))


if __name__ == '__main__':

    try:
        os.system('cls')
    except:
        try:
            os.system('clear')
        except:
                None
    finally:
        print("\n     ###################################################################################################")    
        print('     #  Script started:      '+ str(datetime.now()) + '                                                #')
        print("     ###################################################################################################\n")

        main()

        print("\n\n     ###################################################################################################")    
        print('     #  Script ended:         '+ str(datetime.now()) + '                                              #')
        print("     ###################################################################################################\n\n")
